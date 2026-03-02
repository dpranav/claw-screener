#!/usr/bin/env python3
"""
SAP Procurement Spend Cube - Main Analysis Runner
====================================================
CLI entry point for running comprehensive SAP spend analysis.

Modes:
  full      - Full dataset spend cube + maverick deep dive (8-slide PPTX, DOCX, CSV)
  company   - Single company RSEG analysis (5-slide PPTX, DOCX, CSV)
  fi        - FI document type analysis (5-slide PPTX, DOCX, CSV)

Usage:
  python run_analysis.py --mode full --base-path "C:\\data\\Condor HRP 2025"
  python run_analysis.py --mode company --bukrs 113 --name "Granada"
  python run_analysis.py --mode fi --bukrs 113 --doc-types EC KR KG KE KX --name "Granada"
"""

import argparse
import os
import sys
import time
import pandas as pd
import numpy as np
from pptx.util import Inches, Pt, Emu

# Ensure our package is importable
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from config import (COMPANY_CODES, BUKRS_NAME_MAP, VENDOR_THRESHOLDS,
                    safe_vendor_name, fmt_eur, fmt_pct, fmt_count, get_analysis_date)

from modules.loaders import (load_comprehensive_dataset, load_normalized_dataset,
                              filter_by_company, add_country_column,
                              load_bkpf, load_bseg_for_docs, load_lfa1,
                              enrich_with_vendors, merge_bkpf_into_bseg)

from modules.classify import (prepare_rseg_data, prepare_fi_data,
                               compute_maverick_rseg, compute_maverick_fi,
                               compute_credit_heavy_vendors, compute_credit_heavy_pct_based,
                               compute_single_source)

from modules.metrics import (compute_all_metrics,
                              compute_maverick_by_country, compute_maverick_vendors,
                              compute_maverick_gl, compute_maverick_prctr,
                              compute_maverick_kostl, compute_maverick_by_spend_type,
                              compute_compliance_by_country)

from modules.pptx_builder import (create_presentation,
                                   build_slide_spend_cube_overview,
                                   build_slide_maverick_scorecard,
                                   build_slide_maverick_by_country,
                                   build_slide_maverick_by_dimension,
                                   build_slide_maverick_vendors,
                                   build_slide_maverick_direct_indirect,
                                   build_slide_extreme_risk,
                                   build_slide_remediation,
                                   build_company_executive_overview)

from modules.docx_builder import (create_document, add_cover_page, add_toc,
                                   section_executive_summary, section_sap_tables,
                                   section_join_logic, section_classification,
                                   section_debit_credit, section_country_breakdown,
                                   section_gl_analysis, section_vendor_analysis,
                                   section_maverick, section_savings_opportunities,
                                   section_appendix)

DATE = get_analysis_date()


# ============================================================================
# FULL DATASET ANALYSIS
# ============================================================================

def run_full_analysis(base_path, output_dir):
    """
    Full dataset spend cube with maverick deep dive.
    Generates: 8-slide PPTX, comprehensive DOCX, maverick CSV.
    """
    print("=" * 100)
    print("SAP PROCUREMENT SPEND CUBE - FULL DATASET ANALYSIS")
    print("=" * 100)
    start = time.time()

    # 1. Load data
    print("\n[1/7] Loading data...")
    df = load_comprehensive_dataset(base_path)
    df = add_country_column(df)
    df = prepare_rseg_data(df)

    # 2. Compute ALL metrics
    print("\n[2/7] Computing all metrics...")
    m = compute_all_metrics(df, mode='full')
    by_country = m['by_country']

    # 3. Compute maverick
    print("\n[3/7] Computing maverick spend...")
    mav = compute_maverick_rseg(df)
    mav_mask = mav['primary']['mask']
    ext_mask = mav['critical']['mask']

    mav_country = compute_maverick_by_country(df, mav_mask, by_country)
    mav_vendors = compute_maverick_vendors(df, mav_mask)
    mav_gl = compute_maverick_gl(df, mav_mask)
    mav_prctr = compute_maverick_prctr(df, mav_mask)
    mav_kostl = compute_maverick_kostl(df, mav_mask)
    mav_spend_type = compute_maverick_by_spend_type(df, mav_mask)
    compliance = compute_compliance_by_country(df)

    # Extreme vendors and country
    ext_vendors = compute_maverick_vendors(df, ext_mask)
    ext_country = compute_maverick_by_country(df, ext_mask, by_country)
    ext_gl = compute_maverick_gl(df, ext_mask, top_n=6)

    # Credit-heavy & single-source
    print("  Computing credit-heavy vendors...")
    high_credit = compute_credit_heavy_pct_based(df, min_debit=VENDOR_THRESHOLDS['credit_heavy_min_debit_full'])

    print("  Computing single-source risk...")
    _, ss_count, ss_pct, ms_count, ss_spend, ms_spend = compute_single_source(df)
    ss_info = {'count': ss_count, 'pct': ss_pct, 'spend': ss_spend, 'multi_count': ms_count, 'multi_spend': ms_spend}

    # 4. Print summary
    tg = m['total_gross']
    print(f"\n  FULL DATASET SUMMARY:")
    print(f"  Total Records:     {m['total_records']:,}")
    print(f"  Total Gross:       EUR {tg/1e6:.0f}M")
    print(f"  Direct/Indirect:   {m['direct_pct']:.0f}% / {m['indirect_pct']:.0f}%")
    print(f"  Vendors:           {m.get('total_vendors', 0):,}")
    print(f"  Maverick (NoPR+NoGR): EUR {mav['primary']['spend']/1e6:.0f}M ({mav['primary']['spend']/tg*100:.1f}%)")
    print(f"  Extreme Risk:      EUR {mav['critical']['spend']/1e6:.0f}M")
    print(f"  3-Way Match:       {mav.get('three_way_rate', 0):.1f}%")

    # 5. Save Maverick CSV
    print(f"\n[4/7] Saving Maverick CSV...")
    os.makedirs(output_dir, exist_ok=True)
    mav_df = df[mav_mask].copy()
    mav_df['Maverick_Category'] = 'No PR + No GR'
    mav_df.loc[ext_mask[mav_mask.index][mav_mask], 'Maverick_Category'] = 'No PR + No GR + No Material (EXTREME)'
    csv_mav = os.path.join(output_dir, 'SAP_Maverick_Spend_Full_Dataset.csv')
    mav_df.to_csv(csv_mav, index=False, encoding='utf-8-sig')
    print(f"  Saved: {csv_mav} ({len(mav_df):,} rows)")

    # 6. Create PowerPoint (8 slides)
    print(f"\n[5/7] Creating 8-slide PowerPoint...")
    prs = create_presentation()
    prefix = "CONDOR HRP 2025 - "

    build_slide_spend_cube_overview(prs, m, mav, by_country, m['vendor_spend'], prefix)
    build_slide_maverick_scorecard(prs, m, mav, high_credit, ss_info, prefix)
    build_slide_maverick_by_country(prs, mav, mav_country, ext_country, compliance, prefix)
    build_slide_maverick_by_dimension(prs, mav_gl, mav_prctr, mav_kostl, mav_spend_type, prefix)
    build_slide_maverick_vendors(prs, m, mav, mav_vendors, high_credit, ss_info, mav_spend_type, prefix)
    build_slide_maverick_direct_indirect(prs, m, mav, mav_spend_type, df, prefix)
    build_slide_extreme_risk(prs, mav, ext_vendors, ext_country, ext_gl, df, prefix)
    build_slide_remediation(prs, m, mav, high_credit, ss_info, prefix)

    pptx_out = os.path.join(output_dir, 'SAP_SpendCube_Maverick_DeepDive.pptx')
    prs.save(pptx_out)
    print(f"  Saved: {pptx_out} (8 slides)")

    # 7. Create Word document
    print(f"\n[6/7] Creating comprehensive Word document...")
    doc = create_document()
    add_cover_page(doc, "SAP Procurement Spend Cube", "Full Dataset Analysis",
                   "Maverick Spend + Compliance Deep Dive",
                   f"Generated: {DATE} | {m['total_records']:,} records | {m.get('total_vendors',0):,} vendors")

    # Maverick vendor spend for docx (rename column)
    mav_vendors_doc = mav_vendors.copy()
    if 'Mav_Gross' not in mav_vendors_doc.columns and 'Gross_EUR' in mav_vendors_doc.columns:
        pass  # already has Gross_EUR

    toc_items = [
        '1. Executive Summary', '2. SAP Raw Tables', '3. Join Logic',
        '4. Direct/Indirect Classification', '5. Debit/Credit Analysis',
        '6. Spend by Country', '7. GL Account Analysis',
        '8. Vendor Analysis', '9. Maverick Spend Analysis',
        '10. Savings Opportunities', 'Appendix: Complete Metrics',
    ]
    add_toc(doc, toc_items)

    section_executive_summary(doc, m, "Condor HRP 2025 Full Dataset", "All Companies")
    section_sap_tables(doc)
    section_join_logic(doc, mode='full')
    section_classification(doc, m)
    section_debit_credit(doc, m)
    section_country_breakdown(doc, by_country, tg)
    section_gl_analysis(doc, m)
    section_vendor_analysis(doc, m)
    section_maverick(doc, m, mav, mav_vendors, high_credit, mode='full')
    section_savings_opportunities(doc, m, mav, high_credit, ss_info)
    section_appendix(doc, m, mav, "Condor HRP 2025 Full Dataset", "All")

    docx_out = os.path.join(output_dir, 'SAP_SpendCube_Full_Analysis.docx')
    doc.save(docx_out)
    print(f"  Saved: {docx_out}")

    # Done
    elapsed = time.time() - start
    print(f"\n[7/7] Complete!")
    print(f"{'='*100}")
    print(f"ALL OUTPUTS GENERATED in {elapsed:.1f}s")
    print(f"{'='*100}")
    print(f"  1. {pptx_out} (8 slides)")
    print(f"  2. {docx_out}")
    print(f"  3. {csv_mav} ({len(mav_df):,} maverick rows)")
    print(f"\nKEY FINDINGS:")
    print(f"  Total Spend:       EUR {tg/1e6:.0f}M ({m['total_records']:,} records)")
    print(f"  3-Way Match Rate:  {mav.get('three_way_rate',0):.1f}%")
    print(f"  Maverick Spend:    EUR {mav['primary']['spend']/1e6:.0f}M ({mav['primary']['spend']/tg*100:.1f}%)")
    print(f"  CRITICAL Risk:     EUR {mav['critical']['spend']/1e6:.0f}M")
    print(f"  Credit-Heavy:      {len(high_credit)} vendors")
    print(f"  Single-Source:     {ss_pct:.0f}% of materials")
    print(f"{'='*100}")


# ============================================================================
# SINGLE COMPANY ANALYSIS
# ============================================================================

def run_company_analysis(base_path, output_dir, bukrs, name):
    """Single company RSEG analysis."""
    print("=" * 100)
    print(f"SAP PROCUREMENT ANALYSIS - {name} (BUKRS={bukrs})")
    print("=" * 100)
    start = time.time()

    print("\n[1/5] Loading and filtering data...")
    df = load_comprehensive_dataset(base_path)
    df = filter_by_company(df, bukrs)
    df = add_country_column(df)
    df = prepare_rseg_data(df)

    print("\n[2/5] Computing metrics...")
    m = compute_all_metrics(df, mode='company')

    print("\n[3/5] Computing maverick...")
    mav = compute_maverick_rseg(df)
    mav_mask = mav['primary']['mask']
    mav_vendors = compute_maverick_vendors(df, mav_mask)
    # Rename for compatibility
    mav_vendors_compat = mav_vendors.rename(columns={'Gross_EUR': 'Mav_Gross', 'Lines': 'Mav_Lines'})
    high_credit = compute_credit_heavy_vendors(df, min_debit=VENDOR_THRESHOLDS['credit_heavy_min_debit_company'])
    _, ss_count, ss_pct, _, ss_spend, _ = compute_single_source(df)
    ss_info = {'count': ss_count, 'pct': ss_pct, 'spend': ss_spend}
    mav_spend_type = compute_maverick_by_spend_type(df, mav_mask)

    # Save CSV
    os.makedirs(output_dir, exist_ok=True)
    csv_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_Complete_Dataset.csv')
    df.to_csv(csv_out, index=False, encoding='utf-8-sig')
    print(f"  CSV: {csv_out} ({len(df):,} rows)")

    # PPTX (5 slides)
    print("\n[4/5] Creating PowerPoint...")
    prs = create_presentation()
    build_company_executive_overview(prs, m, name, bukrs)

    from modules.pptx_builder import (add_slide, add_title_bar, add_text, add_kpi, add_table as pptx_add_table,
                                       add_footer, add_insight_box, DARK_BLUE, LIGHT_BLUE, GREEN, SPAIN_RED,
                                       RED, ORANGE, PURPLE, DARK_RED, VN, _f, _p)

    # Slide 2: GL & Cost Center
    s2 = add_slide(prs)
    add_title_bar(s2, f"GL ACCOUNT & COST CENTER - {name.upper()}", "BSEG Enriched", DARK_BLUE)
    gl = m.get('gl_analysis', pd.DataFrame())
    if len(gl) > 0:
        add_text(s2, "Top GL Accounts", Inches(0.5), Inches(1.1), Inches(6), Inches(0.3), font_size=14, bold=True, color=SPAIN_RED)
        gd = [["GL Account", "Total EUR", "Lines"]]
        for _, row in gl.head(12).iterrows():
            gs = str(int(row['GL_Account'])) if pd.notna(row['GL_Account']) else 'N/A'
            gd.append([gs, _f(row['Total_EUR'], 'M1'), f"{int(row['Lines']):,}"])
        pptx_add_table(s2, gd, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5), SPAIN_RED)
    kostl = m.get('kostl_analysis', pd.DataFrame())
    if len(kostl) > 0:
        add_text(s2, "Top Cost Centers", Inches(7.0), Inches(1.1), Inches(6), Inches(0.3), font_size=14, bold=True, color=LIGHT_BLUE)
        kd = [["Cost Center", "Total EUR", "Lines"]]
        for _, row in kostl.iterrows():
            kd.append([str(row['Cost_Center'])[:15], _f(row['Total_EUR'], 'M1'), f"{int(row['Lines']):,}"])
        pptx_add_table(s2, kd, Inches(7.0), Inches(1.4), Inches(6.0), Inches(3.5), LIGHT_BLUE)
    add_footer(s2, f"Source: SAP BSEG | {name} | {DATE}")

    # Slide 3: Vendor
    from modules.pptx_builder import build_slide_maverick_vendors as _bsmv
    # Use the full-dataset vendor slide but with single company data
    s3 = add_slide(prs)
    add_title_bar(s3, f"VENDOR CONCENTRATION - {name.upper()}", "", DARK_BLUE)
    add_kpi(s3, "Total Vendors", f"{m.get('total_vendors',0):,}", Inches(0.5), Inches(1.1), Inches(2.5), Inches(0.8), DARK_BLUE)
    add_kpi(s3, "Top 10 = % Spend", f"{m.get('top10_pct',0):.1f}%", Inches(3.2), Inches(1.1), Inches(2.5), Inches(0.8), SPAIN_RED)
    add_kpi(s3, "Pareto 80%", f"{m.get('pareto_80',0):,} vendors", Inches(5.9), Inches(1.1), Inches(2.5), Inches(0.8), PURPLE)
    vs = m.get('vendor_spend', pd.DataFrame())
    if len(vs) > 0:
        add_text(s3, "Top 15 Vendors by Spend", Inches(0.5), Inches(2.1), Inches(12), Inches(0.3), font_size=14, bold=True, color=DARK_BLUE)
        vd = [["Rank", "Vendor", "Total EUR", "Lines", "% Total"]]
        for i, (_, row) in enumerate(vs.head(15).iterrows(), 1):
            amt = row.get('Total_EUR', row.get('Gross_EUR', 0))
            vd.append([str(i), VN(row['Vendor'], 30), _f(amt, 'M1'), f"{int(row['Lines']):,}", _p(amt, m['total_gross'])])
        pptx_add_table(s3, vd, Inches(0.5), Inches(2.4), Inches(12.3), Inches(4.0), DARK_BLUE)
    add_footer(s3, f"Source: SAP | {name} | {DATE}")

    # Slide 4: Maverick
    from modules.pptx_builder import build_slide_maverick_vendors
    build_slide_maverick_vendors(prs, m, mav, mav_vendors, high_credit, ss_info, mav_spend_type, f"{name} - ")

    # Slide 5: Recommendations
    build_slide_remediation(prs, m, mav, high_credit, ss_info, f"{name} - ")

    pptx_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_Presentation.pptx')
    prs.save(pptx_out)
    print(f"  PPTX: {pptx_out} (5 slides)")

    # DOCX
    print("\n[5/5] Creating Word document...")
    doc = create_document()
    add_cover_page(doc, "SAP Procurement Analysis", f"Company Code: {bukrs}", name,
                   f"Generated: {DATE} | {m['total_records']:,} records")
    section_executive_summary(doc, m, name, bukrs)
    section_sap_tables(doc)
    section_join_logic(doc, mode='rseg')
    section_classification(doc, m)
    section_debit_credit(doc, m)
    section_gl_analysis(doc, m)
    section_vendor_analysis(doc, m)
    section_maverick(doc, m, mav, mav_vendors, high_credit, mode='company')
    section_appendix(doc, m, mav, name, bukrs)

    docx_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_Analysis.docx')
    doc.save(docx_out)
    print(f"  DOCX: {docx_out}")

    elapsed = time.time() - start
    print(f"\nComplete in {elapsed:.1f}s: {len(df):,} records, EUR {m['total_gross']/1e6:.1f}M")


# ============================================================================
# FI DOCUMENT TYPE ANALYSIS
# ============================================================================

def run_fi_analysis(base_path, output_dir, bukrs, name, doc_types):
    """FI document type analysis from BKPF/BSEG raw tables."""
    print("=" * 100)
    print(f"SAP FI DOCUMENT ANALYSIS - {name} (BUKRS={bukrs}) - Doc Types: {', '.join(doc_types)}")
    print("=" * 100)
    start = time.time()

    print("\n[1/5] Loading BKPF, BSEG, LFA1...")
    bkpf = load_bkpf(base_path, bukrs, doc_types)

    target_keys = set(zip(bkpf['BELNR'].astype(str), bkpf['GJAHR'].astype(str)))
    print(f"  Target documents: {len(target_keys):,}")

    bseg = load_bseg_for_docs(base_path, target_keys, bukrs)
    if len(bseg) == 0:
        print("ERROR: No BSEG lines matched. Aborting.")
        return

    bseg = merge_bkpf_into_bseg(bseg, bkpf)

    lfa1 = load_lfa1(base_path)
    bseg = enrich_with_vendors(bseg, lfa1)

    print(f"\n[2/5] Preparing FI data ({len(bseg):,} lines)...")
    df = prepare_fi_data(bseg)

    print("\n[3/5] Computing metrics...")
    m = compute_all_metrics(df, mode='fi')

    mav = compute_maverick_fi(df)
    mav_mask = mav['primary']['mask']
    mav_vendors = compute_maverick_vendors(df, mav_mask)
    high_credit = compute_credit_heavy_vendors(df, min_debit=VENDOR_THRESHOLDS['credit_heavy_min_debit_company'])
    _, ss_count, ss_pct, _, ss_spend, _ = compute_single_source(df)
    ss_info = {'count': ss_count, 'pct': ss_pct, 'spend': ss_spend}
    mav_spend_type = compute_maverick_by_spend_type(df, mav_mask)

    # Save CSV
    os.makedirs(output_dir, exist_ok=True)
    types_str = '_'.join(doc_types)
    csv_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_DocTypes_{types_str}.csv')
    df.to_csv(csv_out, index=False, encoding='utf-8-sig')
    print(f"  CSV: {csv_out} ({len(df):,} rows)")

    # PPTX (5 slides)
    print("\n[4/5] Creating PowerPoint...")
    prs = create_presentation()
    doc_info = ', '.join(doc_types)
    build_company_executive_overview(prs, m, name, bukrs, doc_type_info=doc_info)

    from modules.pptx_builder import (add_slide as _as, add_title_bar as _atb, add_text as _at,
                                       add_kpi as _ak, add_table as _tbl, add_footer as _af,
                                       DARK_BLUE, LIGHT_BLUE, GREEN, SPAIN_RED, RED, ORANGE, PURPLE,
                                       VN, _f, _p)

    # GL & CC slide
    s2 = _as(prs)
    _atb(s2, f"GL ACCOUNT & COST CENTER - {name.upper()}", f"Doc Types: {doc_info}", DARK_BLUE)
    gl = m.get('gl_analysis', pd.DataFrame())
    if len(gl) > 0:
        _at(s2, "Top GL Accounts", Inches(0.5), Inches(1.1), Inches(6), Inches(0.3), font_size=14, bold=True, color=SPAIN_RED)
        gd = [["GL Account", "Total EUR", "Lines"]]
        for _, row in gl.head(12).iterrows():
            gs = str(int(row['GL_Account'])) if pd.notna(row['GL_Account']) else 'N/A'
            gd.append([gs, _f(row['Total_EUR'], 'M1'), f"{int(row['Lines']):,}"])
        _tbl(s2, gd, Inches(0.5), Inches(1.4), Inches(6.0), Inches(3.5), SPAIN_RED)
    _af(s2, f"Source: SAP BSEG | {name} | {DATE}")

    # Vendor slide
    s3 = _as(prs)
    _atb(s3, f"VENDOR CONCENTRATION - {name.upper()}", "", DARK_BLUE)
    _ak(s3, "Total Vendors", f"{m.get('total_vendors',0):,}", Inches(0.5), Inches(1.1), Inches(2.5), Inches(0.8), DARK_BLUE)
    vs = m.get('vendor_spend', pd.DataFrame())
    if len(vs) > 0:
        _at(s3, "Top 15 Vendors", Inches(0.5), Inches(2.1), Inches(12), Inches(0.3), font_size=14, bold=True, color=DARK_BLUE)
        vd = [["Rank", "Vendor", "Total EUR", "Lines", "% Total"]]
        for i, (_, row) in enumerate(vs.head(15).iterrows(), 1):
            amt = row.get('Total_EUR', row.get('Gross_EUR', 0))
            vd.append([str(i), VN(row['Vendor'], 30), _f(amt, 'M1'), f"{int(row['Lines']):,}", _p(amt, m['total_gross'])])
        _tbl(s3, vd, Inches(0.5), Inches(2.4), Inches(12.3), Inches(4.0), DARK_BLUE)
    _af(s3, f"Source: SAP | {name} | {DATE}")

    # Maverick + Recommendations
    build_slide_maverick_vendors(prs, m, mav, mav_vendors, high_credit, ss_info, mav_spend_type, f"{name} FI - ")
    build_slide_remediation(prs, m, mav, high_credit, ss_info, f"{name} FI - ")

    pptx_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_DocTypes_Presentation.pptx')
    prs.save(pptx_out)
    print(f"  PPTX: {pptx_out}")

    # DOCX
    print("\n[5/5] Creating Word document...")
    doc = create_document()
    add_cover_page(doc, "SAP FI Document Analysis", f"Company Code: {bukrs}", f"{name} - {doc_info}",
                   f"Generated: {DATE} | {m['total_records']:,} records")
    section_executive_summary(doc, m, f"{name} FI ({doc_info})", bukrs)
    section_join_logic(doc, mode='fi')
    section_classification(doc, m)
    section_debit_credit(doc, m)
    section_gl_analysis(doc, m)
    section_vendor_analysis(doc, m)
    section_maverick(doc, m, mav, mav_vendors, high_credit, mode='fi')
    section_appendix(doc, m, mav, f"{name} FI ({doc_info})", bukrs)

    docx_out = os.path.join(output_dir, f'SAP_{name}_{bukrs}_DocTypes_Analysis.docx')
    doc.save(docx_out)
    print(f"  DOCX: {docx_out}")

    elapsed = time.time() - start
    print(f"\nComplete in {elapsed:.1f}s: {len(df):,} records, EUR {m['total_gross']/1e6:.1f}M")


# ============================================================================
# MAIN
# ============================================================================

def main():
    parser = argparse.ArgumentParser(description='SAP Procurement Spend Cube Analysis')
    parser.add_argument('--mode', choices=['full', 'company', 'fi'], default='full',
                        help='Analysis mode: full (all companies), company (single), fi (FI doc types)')
    parser.add_argument('--bukrs', type=int, help='Company code for single company or FI analysis')
    parser.add_argument('--name', type=str, help='Entity name (e.g., "Granada")')
    parser.add_argument('--doc-types', nargs='+', help='FI document types (e.g., EC KR KG KE KX)')
    parser.add_argument('--base-path', type=str, default=r'C:\data\Condor HRP 2025',
                        help='Base path to SAP data files')
    parser.add_argument('--output-dir', type=str, help='Output directory (auto-generated if not specified)')

    args = parser.parse_args()

    # Determine output directory
    if args.output_dir:
        out_dir = args.output_dir
    elif args.mode == 'full':
        out_dir = os.path.join(args.base_path, 'final-output-v3', 'SpendCube_Full')
    elif args.mode == 'company':
        name = args.name or COMPANY_CODES.get(args.bukrs, {}).get('name', f'BUKRS_{args.bukrs}')
        out_dir = os.path.join(args.base_path, 'final-output-v3', f'{name}_{args.bukrs}')
    elif args.mode == 'fi':
        name = args.name or COMPANY_CODES.get(args.bukrs, {}).get('name', f'BUKRS_{args.bukrs}')
        out_dir = os.path.join(args.base_path, 'final-output-v3', f'{name}_{args.bukrs}_DocTypes')

    if args.mode == 'full':
        run_full_analysis(args.base_path, out_dir)
    elif args.mode == 'company':
        if not args.bukrs:
            parser.error('--bukrs is required for company mode')
        name = args.name or COMPANY_CODES.get(args.bukrs, {}).get('name', f'BUKRS_{args.bukrs}')
        run_company_analysis(args.base_path, out_dir, args.bukrs, name)
    elif args.mode == 'fi':
        if not args.bukrs or not args.doc_types:
            parser.error('--bukrs and --doc-types are required for fi mode')
        name = args.name or COMPANY_CODES.get(args.bukrs, {}).get('name', f'BUKRS_{args.bukrs}')
        run_fi_analysis(args.base_path, out_dir, args.bukrs, name, args.doc_types)


if __name__ == '__main__':
    main()
