"""
SAP Procurement Spend Cube - PowerPoint Presentation Builder
==============================================================
Creates formatted PPTX decks for both full-dataset (8 slides) and
single-company (5 slides) analyses.
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import (PPTX_COLORS, PPTX_DIMENSIONS, GL_ACCOUNT_DESCRIPTIONS,
                    safe_vendor_name, get_analysis_date)


# ============================================================================
# COLOR HELPERS
# ============================================================================
def rgb(name):
    c = PPTX_COLORS.get(name, (0, 0, 0))
    return RGBColor(c[0], c[1], c[2])

def rgb_tuple(t):
    return RGBColor(t[0], t[1], t[2])

# Pre-built colors
DARK_BLUE = rgb('DARK_BLUE')
LIGHT_BLUE = rgb('LIGHT_BLUE')
GREEN = rgb('ACCENT_GREEN')
SPAIN_RED = rgb('SPAIN_RED')
WHITE = rgb('WHITE')
DARK_GRAY = rgb('DARK_GRAY')
LIGHT_GRAY = rgb('LIGHT_GRAY')
RED = rgb('RED')
DARK_RED = rgb('DARK_RED')
PURPLE = rgb('PURPLE')
ORANGE = rgb('ORANGE')
GOLD = rgb('GOLD')

VN = safe_vendor_name
DATE = get_analysis_date()


# ============================================================================
# PRIMITIVES
# ============================================================================
def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(PPTX_DIMENSIONS['width_inches'])
    prs.slide_height = Inches(PPTX_DIMENSIONS['height_inches'])
    return prs

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_text(slide, text, left, top, width, height, font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
    if color is None: color = DARK_GRAY
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(text); p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return shape

def add_multi_text(slide, lines, left, top, width, height, font_size=9, color=None):
    if color is None: color = DARK_GRAY
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(font_size); p.font.color.rgb = color; p.space_after = Pt(2)
    return shape

def add_title_bar(slide, title, subtitle='', bar_color=None):
    if bar_color is None: bar_color = SPAIN_RED
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    shape.fill.solid(); shape.fill.fore_color.rgb = bar_color; shape.line.fill.background()
    add_text(slide, title, Inches(0.5), Inches(0.15), Inches(9), Inches(0.3), font_size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.5), Inches(9), Inches(0.3), font_size=12, color=WHITE)
    return shape

def add_kpi(slide, title, value, left, top, width=None, height=None, box_color=None):
    if width is None: width = Inches(3.0)
    if height is None: height = Inches(0.8)
    if box_color is None: box_color = DARK_BLUE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = rgb_tuple(PPTX_COLORS['KPI_BG'])
    s.line.color.rgb = box_color; s.line.width = Pt(2)
    add_text(slide, title, left+Inches(0.08), top+Inches(0.05), width-Inches(0.16), Inches(0.25),
             font_size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, value, left+Inches(0.08), top+Inches(0.32), width-Inches(0.16), Inches(0.4),
             font_size=16, bold=True, color=box_color, align=PP_ALIGN.CENTER)

def add_table(slide, data, left, top, width, height, hdr_color=None):
    if hdr_color is None: hdr_color = DARK_BLUE
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height); t = ts.table
    cw = width // cols
    for i in range(cols): t.columns[i].width = cw
    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd):
            cell = t.cell(ri, ci); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.alignment = PP_ALIGN.CENTER
                if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE; cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
                else:
                    p.font.color.rgb = DARK_GRAY
                    if ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    return ts

def add_insight_box(slide, text, left, top, width, height, border_color=None, bg=None):
    if border_color is None: border_color = ORANGE
    if bg is None: bg = PPTX_COLORS['INSIGHT_BG']
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = rgb_tuple(bg); box.line.color.rgb = border_color
    add_text(slide, text, left+Inches(0.2), top+Inches(0.1), width-Inches(0.4), height-Inches(0.2),
             font_size=10, color=DARK_GRAY)

def add_footer(slide, text):
    add_text(slide, text, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3), font_size=9, color=DARK_GRAY)

def _f(v, fmt='M'):
    if fmt == 'M': return f"EUR {v/1e6:.0f}M"
    elif fmt == 'M1': return f"EUR {v/1e6:.1f}M"
    elif fmt == 'M2': return f"EUR {v/1e6:.2f}M"
    elif fmt == 'K': return f"EUR {v/1e3:.0f}K"
    return f"EUR {v:,.0f}"

def _p(v, t): return f"{v/t*100:.1f}%" if t > 0 else "0%"


# ============================================================================
# SLIDE 1: SPEND CUBE OVERVIEW (Full Dataset)
# ============================================================================
def build_slide_spend_cube_overview(prs, m, mav, by_country, vendor_spend, title_prefix=""):
    """Full dataset spend cube overview with country and vendor breakdown."""
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}SPEND CUBE OVERVIEW", "Full Dataset | BSEG + Vendor Enhanced", DARK_BLUE)

    tg = m['total_gross']; tn = m['total_net']; tr = m['total_records']; tv = m.get('total_vendors', 0)
    ds = m['direct_spend']; ins = m['indirect_spend']

    add_kpi(s, "Total Gross Spend", _f(tg), Inches(0.5), Inches(1.1), Inches(2.5), Inches(0.8), DARK_BLUE)
    add_kpi(s, "Total Net Spend", _f(tn), Inches(3.2), Inches(1.1), Inches(2.5), Inches(0.8), DARK_BLUE)
    add_kpi(s, "Invoice Lines", f"{tr:,}", Inches(5.9), Inches(1.1), Inches(2.2), Inches(0.8), DARK_BLUE)
    add_kpi(s, "Vendors", f"{tv:,}", Inches(8.3), Inches(1.1), Inches(2.2), Inches(0.8), GREEN)
    add_kpi(s, "Direct / Indirect", f"{_p(ds,tg)} / {_p(ins,tg)}", Inches(10.7), Inches(1.1), Inches(2.3), Inches(0.8), ORANGE)

    # Country table
    add_text(s, "Spend by Country", Inches(0.5), Inches(2.1), Inches(6), Inches(0.3), font_size=13, bold=True, color=DARK_BLUE)
    cd = [["Country", "Gross EUR", "Net EUR", "Lines", "% Total"]]
    for _, r in by_country.iterrows():
        cd.append([str(r['Country'])[:20], _f(r['Gross_EUR']), _f(r['Net_EUR']), f"{int(r['Lines']):,}", _p(r['Gross_EUR'], tg)])
    add_table(s, cd, Inches(0.5), Inches(2.4), Inches(6.3), Inches(3.2), DARK_BLUE)

    # Top vendors
    add_text(s, "Top 10 Vendors", Inches(7.2), Inches(2.1), Inches(5.8), Inches(0.3), font_size=13, bold=True, color=GREEN)
    vd = [["Rank", "Vendor", "Gross EUR", "% Total"]]
    for i, (_, r) in enumerate(vendor_spend.head(10).iterrows(), 1):
        vd.append([str(i), VN(r['Vendor']), _f(r.get('Gross_EUR', r.get('Total_EUR', 0))), _p(r.get('Gross_EUR', r.get('Total_EUR', 0)), tg)])
    add_table(s, vd, Inches(7.2), Inches(2.4), Inches(5.8), Inches(3.2), GREEN)

    # BSEG coverage box
    bseg_txt = f"BSEG: GL {m['gl_pct']:.0f}% | PC {m['prctr_pct']:.0f}% | CC {m['kostl_pct']:.0f}% | Vendors: {tv:,} ({m.get('vendor_coverage',0):.0f}%) | Top10={m.get('top10_pct',0):.0f}% | Pareto 80%={m.get('pareto_80',0):,}"
    add_insight_box(s, bseg_txt, Inches(0.5), Inches(5.8), Inches(6.3), Inches(1.0), LIGHT_BLUE, PPTX_COLORS['BSEG_BG'])

    # Key insight
    top_country = by_country.iloc[0]['Country'] if len(by_country) > 0 else 'N/A'
    top_pct = by_country.iloc[0]['Gross_EUR'] / tg * 100 if len(by_country) > 0 and tg > 0 else 0
    add_insight_box(s, f"{top_country} dominates at {top_pct:.0f}% of spend. {tv:,} vendors (Top10={m.get('top10_pct',0):.0f}%). {_p(ins,tg)} indirect spend opportunity.",
                    Inches(7.2), Inches(5.8), Inches(5.8), Inches(1.0), ORANGE)

    add_footer(s, f"Source: SAP Final Comprehensive with BSEG + Vendor | {tr:,} records | {DATE}")
    return s


# ============================================================================
# SLIDE 2: MAVERICK SCORECARD & THREE-WAY MATCH
# ============================================================================
def build_slide_maverick_scorecard(prs, m, mav, high_credit, ss_info, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK SPEND SCORECARD & THREE-WAY MATCH", "Full Dataset", DARK_RED)

    tg = mav['total_gross']; tr = mav['total_records']

    add_kpi(s, "CRITICAL Risk", f"{_f(mav['critical']['spend'])} ({mav['critical']['lines']:,})", Inches(0.5), Inches(1.1), Inches(3.0), Inches(0.8), DARK_RED)
    add_kpi(s, "HIGH: No PR+No GR", f"{_f(mav['primary']['spend'])} ({mav['primary']['lines']:,})", Inches(3.7), Inches(1.1), Inches(3.0), Inches(0.8), RED)
    add_kpi(s, "3-Way Match Rate", f"{mav.get('three_way_rate', 0):.1f}%", Inches(6.9), Inches(1.1), Inches(3.0), Inches(0.8), GREEN)
    tw_spend = mav.get('three_way_match', {}).get('spend', 0)
    add_kpi(s, "Compliant Spend", _f(tw_spend), Inches(10.1), Inches(1.1), Inches(3.0), Inches(0.8), GREEN)

    # Scorecard
    add_text(s, "Maverick Risk Scorecard", Inches(0.5), Inches(2.1), Inches(7), Inches(0.3), font_size=14, bold=True, color=DARK_RED)
    no_pr = mav.get('medium_no_pr', mav.get('primary', {}))
    no_gr = mav.get('medium_no_gr', {'lines': 0, 'spend': 0})
    sc = [["Risk", "Category", "Lines", "% Lines", "Spend EUR", "% Spend"],
          ["CRITICAL", "No PR+GR+Mat", f"{mav['critical']['lines']:,}", _p(mav['critical']['lines'], tr), _f(mav['critical']['spend']), _p(mav['critical']['spend'], tg)],
          ["HIGH", "No PR + No GR", f"{mav['primary']['lines']:,}", _p(mav['primary']['lines'], tr), _f(mav['primary']['spend']), _p(mav['primary']['spend'], tg)],
          ["MEDIUM", "No PR", f"{no_pr.get('lines',0):,}", _p(no_pr.get('lines',0), tr), _f(no_pr.get('spend',0)), _p(no_pr.get('spend',0), tg)],
          ["MEDIUM", "No GR", f"{no_gr.get('lines',0):,}", _p(no_gr.get('lines',0), tr), _f(no_gr.get('spend',0)), _p(no_gr.get('spend',0), tg)],
          ["LOW", "No Material", f"{m.get('no_material_lines',0):,}", _p(m.get('no_material_lines',0), tr), _f(m.get('no_material_spend',0)), _p(m.get('no_material_spend',0), tg)],
          ["MEDIUM", "Credit-Heavy (>20%)", f"{len(high_credit)}", "", _f(high_credit['Credit_EUR'].sum()) if len(high_credit)>0 and 'Credit_EUR' in high_credit.columns else _f(high_credit['Credit_Total'].sum()) if len(high_credit)>0 else "EUR 0M", ""]]
    add_table(s, sc, Inches(0.5), Inches(2.4), Inches(7.5), Inches(2.5), DARK_RED)

    # Three-Way Match
    add_text(s, "Three-Way Match Analysis", Inches(8.3), Inches(2.1), Inches(4.7), Inches(0.3), font_size=14, bold=True, color=DARK_BLUE)
    twm = mav.get('three_way_match', {'lines': 0, 'spend': 0})
    po_gr = mav.get('po_gr', {'lines': 0, 'spend': 0})
    po_pr = mav.get('po_pr', {'lines': 0, 'spend': 0})
    po_only = mav.get('po_only', {'lines': 0, 'spend': 0})
    tm = [["Match Level", "Lines", "%", "Spend EUR"],
          ["Full (PO+PR+GR)", f"{twm['lines']:,}", _p(twm['lines'], tr), _f(twm['spend'])],
          ["PO+GR (No PR)", f"{po_gr['lines']:,}", _p(po_gr['lines'], tr), _f(po_gr['spend'])],
          ["PO+PR (No GR)", f"{po_pr['lines']:,}", _p(po_pr['lines'], tr), _f(po_pr['spend'])],
          ["PO Only MAVERICK", f"{po_only['lines']:,}", _p(po_only['lines'], tr), _f(po_only['spend'])]]
    add_table(s, tm, Inches(8.3), Inches(2.4), Inches(4.7), Inches(2.5), DARK_BLUE)

    # Single source
    ss_count, ss_pct, ss_spend = ss_info.get('count', 0), ss_info.get('pct', 0), ss_info.get('spend', 0)
    add_text(s, "Single-Source & Concentration Risk", Inches(6.0), Inches(5.1), Inches(6), Inches(0.3), font_size=13, bold=True, color=PURPLE)
    ssr = [["Metric", "Value"],
           ["Single-Source Materials", f"{ss_count:,} ({ss_pct:.0f}%)"],
           ["Single-Source Spend", _f(ss_spend)],
           ["Credit-Heavy Vendors", f"{len(high_credit)} vendors"]]
    add_table(s, ssr, Inches(6.0), Inches(5.4), Inches(6.0), Inches(1.0), PURPLE)

    # Insight
    mav_pct = mav['primary']['spend'] / tg * 100 if tg > 0 else 0
    add_insight_box(s, f"KEY: {mav.get('three_way_rate',0):.1f}% fully compliant. {_f(mav['primary']['spend'])} ({mav_pct:.1f}%) maverick across {mav['primary']['lines']:,} transactions. {_f(mav['critical']['spend'])} CRITICAL with zero controls.",
                    Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), DARK_RED, PPTX_COLORS['KEY_FIND_BG'])

    add_footer(s, f"Source: SAP BSEG Enhanced | Maverick = No PR + No GR | {DATE}")
    return s


# ============================================================================
# SLIDE 3: MAVERICK BY COUNTRY
# ============================================================================
def build_slide_maverick_by_country(prs, mav, mav_country, ext_country, compliance, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK SPEND BY COMPANY / COUNTRY", "", DARK_RED)

    add_text(s, "Maverick Spend by Country (No PR + No GR)", Inches(0.5), Inches(1.1), Inches(7), Inches(0.3), font_size=14, bold=True, color=DARK_RED)
    mc = [["Country", "Total Spend", "Maverick Spend", "Mav Lines", "Mav %"]]
    for _, r in mav_country.iterrows():
        mc.append([str(r['Country'])[:20], _f(r['Gross_EUR']), _f(r['Mav_Gross'], 'M1'), f"{int(r['Mav_Lines']):,}", f"{r['Mav_Pct']:.1f}%"])
    add_table(s, mc, Inches(0.5), Inches(1.4), Inches(7.0), Inches(3.5), DARK_RED)

    add_text(s, "CRITICAL Risk by Country", Inches(7.8), Inches(1.1), Inches(5.2), Inches(0.3), font_size=14, bold=True, color=DARK_RED)
    ec = [["Country", "Extreme Spend", "Lines"]]
    for _, r in ext_country.iterrows():
        ec.append([str(r['Country'])[:20], _f(r.get('Ext_Gross', r.get('Gross_EUR', 0)), 'M1'), f"{int(r.get('Ext_Lines', r.get('Lines', 0))):,}"])
    add_table(s, ec, Inches(7.8), Inches(1.4), Inches(5.2), Inches(3.5), DARK_RED)

    # Compliance by country
    if len(compliance) > 0:
        add_text(s, "Compliance Rate by Country", Inches(0.5), Inches(5.2), Inches(12), Inches(0.3), font_size=14, bold=True, color=GREEN)
        cd = [["Country", "Total Lines", "3-Way Match", "Match %", "Maverick", "Mav %"]]
        for _, r in compliance.iterrows():
            cd.append([str(r['Country'])[:20], f"{int(r['Total_Lines']):,}", f"{int(r['Three_Way_Match']):,}",
                        f"{r['Match_Pct']:.1f}%", f"{int(r['Maverick']):,}", f"{r['Mav_Pct']:.1f}%"])
        add_table(s, cd, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), GREEN)

    add_footer(s, f"Source: SAP BSEG Enhanced | {DATE}")
    return s


# ============================================================================
# SLIDE 4: MAVERICK BY GL / CC / PC
# ============================================================================
def build_slide_maverick_by_dimension(prs, mav_gl, mav_prctr, mav_kostl, mav_spend_type, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK BY GL ACCOUNT, COST CENTER & PROFIT CENTER", "", DARK_RED)

    # Maverick GL
    add_text(s, "Maverick by GL Account", Inches(0.5), Inches(1.1), Inches(4.5), Inches(0.3), font_size=13, bold=True, color=DARK_RED)
    gd = [["GL Account", "Description", "Spend", "Lines"]]
    for _, r in mav_gl.iterrows():
        gs = str(int(r['GL_Account'])) if pd.notna(r['GL_Account']) else 'N/A'
        desc = GL_ACCOUNT_DESCRIPTIONS.get(int(float(r['GL_Account'])), '')[:18] if pd.notna(r['GL_Account']) else ''
        gd.append([gs, desc, _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
    add_table(s, gd, Inches(0.5), Inches(1.4), Inches(4.3), Inches(3.5), DARK_RED)

    # Maverick Profit Center
    add_text(s, "Maverick by Profit Center", Inches(5.0), Inches(1.1), Inches(4), Inches(0.3), font_size=13, bold=True, color=PURPLE)
    pd2 = [["Profit Center", "Spend", "Lines"]]
    for _, r in mav_prctr.iterrows():
        pd2.append([str(r['Profit_Center'])[:15], _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
    add_table(s, pd2, Inches(5.0), Inches(1.4), Inches(4.0), Inches(3.5), PURPLE)

    # Maverick Cost Center
    add_text(s, "Maverick by Cost Center", Inches(9.2), Inches(1.1), Inches(3.8), Inches(0.3), font_size=13, bold=True, color=LIGHT_BLUE)
    cd = [["Cost Center", "Spend", "Lines"]]
    for _, r in mav_kostl.iterrows():
        cd.append([str(r['Cost_Center'])[:14], _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
    add_table(s, cd, Inches(9.2), Inches(1.4), Inches(3.8), Inches(3.5), LIGHT_BLUE)

    # Maverick indirect reasons
    mav_ir = mav_spend_type.get('mav_indirect_reasons', pd.DataFrame())
    if len(mav_ir) > 0:
        add_text(s, "Maverick Indirect - Classification Reasons", Inches(0.5), Inches(5.1), Inches(6), Inches(0.3), font_size=13, bold=True, color=ORANGE)
        ir = [["Reason", "Spend", "Lines"]]
        for _, r in mav_ir.head(5).iterrows():
            reason = str(r['Spend_Type_Reason']).replace("Material number has ", "Mat ").replace(" digits", "d").replace("(7+ = Indirect)", "").strip()[:35]
            ir.append([reason, _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
        add_table(s, ir, Inches(0.5), Inches(5.4), Inches(6.5), Inches(1.5), ORANGE)

    # GL patterns
    add_text(s, "GL Pattern Insight", Inches(7.3), Inches(5.1), Inches(5.7), Inches(0.3), font_size=13, bold=True, color=DARK_RED)
    add_insight_box(s, "", Inches(7.3), Inches(5.4), Inches(5.7), Inches(1.5), DARK_RED, PPTX_COLORS['REC_BG'])
    add_multi_text(s, [
        "29xxxxx (GR/IR Clearing) = Largest maverick GL",
        "  Invoices posted to clearing without proper GR",
        "44xxxxx (Materials/Freight) = Direct cost leakage",
        "  Material purchases bypassing 3-way match",
        "67xxxxx (External Services) = Service maverick",
        "  Third-party services without proper controls",
    ], Inches(7.5), Inches(5.5), Inches(5.3), Inches(1.3), font_size=9, color=DARK_GRAY)

    add_footer(s, f"Source: SAP BSEG Enhanced | {DATE}")
    return s


# ============================================================================
# SLIDE 5: MAVERICK VENDOR DEEP DIVE
# ============================================================================
def build_slide_maverick_vendors(prs, m, mav, mav_vendors, high_credit, ss_info, mav_spend_type, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK VENDOR DEEP DIVE & CREDIT RISK", "", RED)

    add_kpi(s, "Maverick Vendors", f"{len(mav_vendors):,}", Inches(0.5), Inches(1.1), Inches(2.5), Inches(0.8), RED)
    add_kpi(s, "Credit-Heavy (>20%)", f"{len(high_credit):,}", Inches(3.2), Inches(1.1), Inches(2.5), Inches(0.8), ORANGE)
    ss_count = ss_info.get('count', 0); ss_pct = ss_info.get('pct', 0)
    add_kpi(s, "Single-Source", f"{ss_count:,} ({ss_pct:.0f}%)", Inches(5.9), Inches(1.1), Inches(2.5), Inches(0.8), PURPLE)
    add_kpi(s, "Mav. Direct", _f(mav_spend_type.get('mav_direct_spend', 0)), Inches(8.6), Inches(1.1), Inches(2.1), Inches(0.8), GREEN)
    add_kpi(s, "Mav. Indirect", _f(mav_spend_type.get('mav_indirect_spend', 0)), Inches(10.9), Inches(1.1), Inches(2.1), Inches(0.8), ORANGE)

    # Top maverick vendors
    mav_total = mav['primary']['spend']
    add_text(s, "Top 15 Maverick Vendors (No PR + No GR)", Inches(0.5), Inches(2.1), Inches(7), Inches(0.3), font_size=13, bold=True, color=RED)
    mv = [["Rank", "Vendor", "Mav Spend", "Lines", "POs", "% Mav"]]
    for i, (_, r) in enumerate(mav_vendors.head(15).iterrows(), 1):
        mv.append([str(i), VN(r['Vendor']), _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}",
                    f"{int(r.get('POs', 0)):,}", _p(r['Gross_EUR'], mav_total)])
    add_table(s, mv, Inches(0.5), Inches(2.4), Inches(7.0), Inches(4.2), RED)

    # Credit-heavy vendors
    add_text(s, "Credit-Heavy Vendors (>20% Ratio)", Inches(7.8), Inches(2.1), Inches(5.2), Inches(0.3), font_size=13, bold=True, color=ORANGE)
    cr = [["Vendor", "Debit", "Credit", "Ratio"]]
    debit_col = 'Debit_EUR' if 'Debit_EUR' in high_credit.columns else 'Debit_Total'
    credit_col = 'Credit_EUR' if 'Credit_EUR' in high_credit.columns else 'Credit_Total'
    for _, r in high_credit.head(12).iterrows():
        cr.append([VN(r['Vendor'], 22), _f(r[debit_col], 'K'), _f(r[credit_col], 'K'), f"{r['Credit_Ratio']:.0f}%"])
    add_table(s, cr, Inches(7.8), Inches(2.4), Inches(5.2), Inches(4.2), ORANGE)

    add_footer(s, f"Source: SAP BSEG Enhanced + Vendor | {DATE}")
    return s


# ============================================================================
# SLIDE 6: MAVERICK DIRECT VS INDIRECT
# ============================================================================
def build_slide_maverick_direct_indirect(prs, m, mav, mav_spend_type, df, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK SPEND - DIRECT VS INDIRECT BREAKDOWN", "", ORANGE)

    mav_total = mav['primary']['spend']
    md = mav_spend_type.get('mav_direct_spend', 0)
    mi = mav_spend_type.get('mav_indirect_spend', 0)

    add_kpi(s, "Total Maverick", _f(mav_total), Inches(0.5), Inches(1.1), Inches(3.0), Inches(0.8), DARK_RED)
    add_kpi(s, "Maverick DIRECT", f"{_f(md)} ({_p(md, mav_total)})" if mav_total > 0 else "EUR 0M", Inches(3.7), Inches(1.1), Inches(3.0), Inches(0.8), GREEN)
    add_kpi(s, "Maverick INDIRECT", f"{_f(mi)} ({_p(mi, mav_total)})" if mav_total > 0 else "EUR 0M", Inches(6.9), Inches(1.1), Inches(3.0), Inches(0.8), ORANGE)
    add_kpi(s, "CRITICAL Risk", _f(mav['critical']['spend']), Inches(10.1), Inches(1.1), Inches(3.0), Inches(0.8), DARK_RED)

    # Direct vs Indirect overall vs maverick comparison
    ds = m['direct_spend']; ins = m['indirect_spend']; tg = m['total_gross']
    add_text(s, "Direct/Indirect: Overall vs Maverick", Inches(0.5), Inches(2.1), Inches(7), Inches(0.3), font_size=13, bold=True, color=DARK_BLUE)
    cmp = [["Metric", "Overall", "Maverick Only", "Delta"],
           ["Direct %", _p(ds, tg), _p(md, mav_total) if mav_total > 0 else "N/A", f"{md/mav_total*100 - ds/tg*100:+.1f}pp" if mav_total > 0 and tg > 0 else ""],
           ["Indirect %", _p(ins, tg), _p(mi, mav_total) if mav_total > 0 else "N/A", f"{mi/mav_total*100 - ins/tg*100:+.1f}pp" if mav_total > 0 and tg > 0 else ""],
           ["Direct Spend", _f(ds), _f(md), f"{_p(md, ds)} maverick" if ds > 0 else ""],
           ["Indirect Spend", _f(ins), _f(mi), f"{_p(mi, ins)} maverick" if ins > 0 else ""]]
    add_table(s, cmp, Inches(0.5), Inches(2.4), Inches(7.0), Inches(1.5), DARK_BLUE)

    # Maverick indirect reasons
    mav_ir = mav_spend_type.get('mav_indirect_reasons', pd.DataFrame())
    if len(mav_ir) > 0:
        add_text(s, "Maverick Indirect Reasons", Inches(7.8), Inches(2.1), Inches(5), Inches(0.3), font_size=13, bold=True, color=ORANGE)
        irr = [["Reason", "Spend", "Lines"]]
        for _, r in mav_ir.head(4).iterrows():
            reason = str(r['Spend_Type_Reason']).replace("Material number has ", "Mat ").replace(" digits","d").replace("(7+ = Indirect)","").strip()[:30]
            irr.append([reason, _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
        add_table(s, irr, Inches(7.8), Inches(2.4), Inches(5.2), Inches(1.5), ORANGE)

    # Top maverick direct vendors
    mav_mask = ~df['Has_PR'] & ~df['Has_GR']
    mav_d = df[mav_mask & (df['Spend_Type'] == 'DIRECT') & df['NAME1'].notna()]
    mav_i = df[mav_mask & (df['Spend_Type'] == 'INDIRECT') & df['NAME1'].notna()]

    if len(mav_d) > 0:
        add_text(s, "Top Maverick DIRECT Vendors", Inches(0.5), Inches(4.2), Inches(6), Inches(0.3), font_size=13, bold=True, color=GREEN)
        dv_agg = mav_d.groupby('NAME1').agg(Gross=('Gross_Amount', 'sum'), Lines=('BELNR', 'count')).reset_index().sort_values('Gross', ascending=False)
        dv = [["Rank", "Vendor", "Mav Direct Spend", "Lines"]]
        for i, (_, r) in enumerate(dv_agg.head(8).iterrows(), 1):
            dv.append([str(i), VN(r['NAME1']), _f(r['Gross'], 'M1'), f"{int(r['Lines']):,}"])
        add_table(s, dv, Inches(0.5), Inches(4.5), Inches(6.3), Inches(2.3), GREEN)

    if len(mav_i) > 0:
        add_text(s, "Top Maverick INDIRECT Vendors", Inches(7.0), Inches(4.2), Inches(6), Inches(0.3), font_size=13, bold=True, color=ORANGE)
        iv_agg = mav_i.groupby('NAME1').agg(Gross=('Gross_Amount', 'sum'), Lines=('BELNR', 'count')).reset_index().sort_values('Gross', ascending=False)
        iv = [["Rank", "Vendor", "Mav Indirect Spend", "Lines"]]
        for i, (_, r) in enumerate(iv_agg.head(8).iterrows(), 1):
            iv.append([str(i), VN(r['NAME1']), _f(r['Gross'], 'M1'), f"{int(r['Lines']):,}"])
        add_table(s, iv, Inches(7.0), Inches(4.5), Inches(6.0), Inches(2.3), ORANGE)

    add_footer(s, f"Source: SAP BSEG Enhanced | {DATE}")
    return s


# ============================================================================
# SLIDE 7: EXTREME RISK DEEP DIVE
# ============================================================================
def build_slide_extreme_risk(prs, mav, ext_vendors, ext_country, ext_gl, df, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}EXTREME RISK: NO PR + NO GR + NO MATERIAL", "Zero Procurement Controls", DARK_RED)

    ext_spend = mav['critical']['spend']; ext_count = mav['critical']['lines']
    tg = mav['total_gross']; mav_spend = mav['primary']['spend']

    add_kpi(s, "EXTREME Lines", f"{ext_count:,}", Inches(0.5), Inches(1.1), Inches(2.5), Inches(0.8), DARK_RED)
    add_kpi(s, "EXTREME Spend", _f(ext_spend), Inches(3.2), Inches(1.1), Inches(2.5), Inches(0.8), DARK_RED)
    add_kpi(s, "% of Total Spend", _p(ext_spend, tg), Inches(5.9), Inches(1.1), Inches(2.5), Inches(0.8), DARK_RED)
    add_kpi(s, "% of Maverick", _p(ext_spend, mav_spend) if mav_spend > 0 else "N/A", Inches(8.6), Inches(1.1), Inches(2.5), Inches(0.8), RED)
    add_kpi(s, "Ext. Vendors", f"{len(ext_vendors):,}", Inches(11.3), Inches(1.1), Inches(1.7), Inches(0.8), RED)

    # Extreme vendors
    add_text(s, "Top 15 EXTREME Risk Vendors", Inches(0.5), Inches(2.1), Inches(7), Inches(0.3), font_size=13, bold=True, color=DARK_RED)
    ev = [["Rank", "Vendor", "Extreme Spend", "Lines", "% Extreme"]]
    for i, (_, r) in enumerate(ext_vendors.head(15).iterrows(), 1):
        ev.append([str(i), VN(r['Vendor']), _f(r['Gross_EUR'], 'M2'), f"{int(r['Lines']):,}", _p(r['Gross_EUR'], ext_spend)])
    add_table(s, ev, Inches(0.5), Inches(2.4), Inches(7.0), Inches(4.2), DARK_RED)

    # Extreme by country
    add_text(s, "EXTREME Risk by Country", Inches(7.8), Inches(2.1), Inches(5.2), Inches(0.3), font_size=13, bold=True, color=DARK_RED)
    ecd = [["Country", "Extreme Spend", "Lines", "% Extreme"]]
    for _, r in ext_country.iterrows():
        g = r.get('Ext_Gross', r.get('Gross_EUR', 0)); l = r.get('Ext_Lines', r.get('Lines', 0))
        ecd.append([str(r['Country'])[:20], _f(g, 'M1'), f"{int(l):,}", _p(g, ext_spend)])
    add_table(s, ecd, Inches(7.8), Inches(2.4), Inches(5.2), Inches(2.0), DARK_RED)

    # Extreme GL
    if len(ext_gl) > 0:
        add_text(s, "EXTREME by GL Account", Inches(7.8), Inches(4.6), Inches(5.2), Inches(0.3), font_size=13, bold=True, color=DARK_RED)
        egd = [["GL Account", "Spend", "Lines"]]
        for _, r in ext_gl.iterrows():
            gs = str(int(r['GL_Account'])) if pd.notna(r['GL_Account']) else 'N/A'
            egd.append([gs, _f(r['Gross_EUR'], 'M1'), f"{int(r['Lines']):,}"])
        add_table(s, egd, Inches(7.8), Inches(4.9), Inches(5.2), Inches(1.7), DARK_RED)

    add_footer(s, f"Source: SAP BSEG Enhanced | EXTREME = No PR + No GR + No Material | {DATE}")
    return s


# ============================================================================
# SLIDE 8: REMEDIATION ROADMAP
# ============================================================================
def build_slide_remediation(prs, m, mav, high_credit, ss_info, title_prefix=""):
    s = add_slide(prs)
    add_title_bar(s, f"{title_prefix}MAVERICK REMEDIATION ROADMAP & EXECUTIVE SUMMARY", "", DARK_BLUE)

    tg = m['total_gross']; tr = m['total_records']; tv = m.get('total_vendors', 0)
    mav_spend = mav['primary']['spend']; ext_spend = mav['critical']['spend']
    tw_spend = mav.get('three_way_match', {}).get('spend', 0)

    # Summary table
    add_text(s, "Executive Summary", Inches(0.5), Inches(1.1), Inches(12), Inches(0.3), font_size=14, bold=True, color=DARK_BLUE)
    hc_credits = high_credit['Credit_EUR'].sum() if 'Credit_EUR' in high_credit.columns else high_credit['Credit_Total'].sum() if len(high_credit)>0 else 0
    ss_count = ss_info.get('count', 0); ss_pct = ss_info.get('pct', 0); ss_spend = ss_info.get('spend', 0)
    es = [["Metric", "Value", "Context"],
          ["Total Spend", _f(tg), f"{tr:,} invoice lines, {tv:,} vendors"],
          ["3-Way Match Rate", f"{mav.get('three_way_rate',0):.1f}%", f"{_f(tw_spend)} compliant spend"],
          ["Total Maverick", f"{_f(mav_spend)} ({_p(mav_spend,tg)})", f"{mav['primary']['lines']:,} transactions, No PR + No GR"],
          ["CRITICAL Risk", f"{_f(ext_spend)} ({_p(ext_spend,tg)})", f"{mav['critical']['lines']:,} txns with ZERO controls"],
          ["Credit-Heavy Vendors", f"{len(high_credit):,} vendors", _f(hc_credits, 'M1') + " in credits" if hc_credits > 0 else ""],
          ["Single-Source", f"{ss_pct:.0f}% of materials", f"{_f(ss_spend)} supply chain risk"]]
    add_table(s, es, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.0), DARK_BLUE)

    # Remediation priorities
    no_pr_spend = mav.get('medium_no_pr', mav['primary']).get('spend', 0)
    no_gr_spend = mav.get('medium_no_gr', {'spend': 0}).get('spend', 0)
    pareto = m.get('pareto_80', 0)
    no_mat_lines = m.get('no_material_lines', 0)

    add_text(s, "Remediation Priorities", Inches(0.5), Inches(3.6), Inches(6), Inches(0.3), font_size=14, bold=True, color=DARK_RED)
    rbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.9), Inches(6.3), Inches(3.0))
    rbox.fill.solid(); rbox.fill.fore_color.rgb = rgb_tuple(PPTX_COLORS['REC_BG']); rbox.line.color.rgb = DARK_RED; rbox.line.width = Pt(2)

    recs = [
        f"1. CRITICAL: Enforce PR for all POs - {_f(no_pr_spend)} bypassing approval",
        f"2. CRITICAL: Mandate GR before payment - {_f(no_gr_spend)} without receipt",
        f"3. HIGH: Auto-block PO creation without approved PR in SAP",
        f"4. HIGH: Implement ERS for top {pareto} vendors",
        f"5. MEDIUM: Assign material numbers to {no_mat_lines:,} service lines",
        f"6. MEDIUM: Review {len(high_credit)} credit-heavy vendors for overbilling",
        f"7. LOW: Dual-source top {ss_count:,} single-source materials",
    ]
    rtxt = s.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(5.9), Inches(2.8))
    tf = rtxt.text_frame; tf.word_wrap = True
    for i, rec in enumerate(recs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = rec; p.font.size = Pt(9); p.font.color.rgb = DARK_GRAY; p.space_after = Pt(3)
        if 'CRITICAL' in rec: p.font.color.rgb = DARK_RED; p.font.bold = True
        elif 'HIGH' in rec: p.font.color.rgb = RED

    # Savings opportunity
    ins = m['indirect_spend']
    add_text(s, "Potential Savings & Impact", Inches(7.0), Inches(3.6), Inches(6), Inches(0.3), font_size=14, bold=True, color=GREEN)
    sbox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(3.9), Inches(6.0), Inches(3.0))
    sbox.fill.solid(); sbox.fill.fore_color.rgb = rgb_tuple(PPTX_COLORS['SAVINGS_BG']); sbox.line.color.rgb = GREEN; sbox.line.width = Pt(2)

    sv = [["Opportunity", "Addressable", "Est. Savings"],
          ["Maverick Compliance", _f(mav_spend), _f(mav_spend * 0.05) + " (5% of maverick)"],
          ["Vendor Consolidation", _f(tg), _f(tg * 0.03) + " (3% of total)"],
          ["Credit Recovery", _f(hc_credits) if hc_credits > 0 else "EUR 0M", "Investigation needed"],
          ["Dual-Source Leverage", _f(ss_spend), _f(ss_spend * 0.02) + " (2% via competition)"],
          ["Indirect Optimization", _f(ins), _f(ins * 0.07) + " (7% category mgmt)"]]
    add_table(s, sv, Inches(7.2), Inches(4.1), Inches(5.6), Inches(1.8), GREEN)

    add_footer(s, f"Source: SAP BSEG Enhanced + Vendor | Full Dataset | {DATE}")
    return s


# ============================================================================
# SINGLE-COMPANY SLIDES (reused from existing skill, adapted)
# ============================================================================
def build_company_executive_overview(prs, metrics, entity_name, entity_code, doc_type_info=None):
    """Build Slide 1 for single company analysis."""
    s = add_slide(prs)
    m = metrics
    title = f"{entity_name} - PROCUREMENT SPEND OVERVIEW"
    subtitle = f"Company Code: {entity_code} | FY 2025"
    if doc_type_info:
        title = f"{entity_name} - FI DOC TYPES: {doc_type_info}"
        subtitle = f"BKPF + BSEG Analysis | Company Code {entity_code} | FY 2025"

    add_title_bar(s, title, subtitle)

    add_kpi(s, "Total Gross Spend", f"EUR {m['total_gross']/1e6:.1f}M", Inches(0.5), Inches(1.1), box_color=SPAIN_RED)
    add_kpi(s, "Net Spend", f"EUR {m['total_net']/1e6:.1f}M", Inches(3.7), Inches(1.1), box_color=DARK_BLUE)
    add_kpi(s, "Invoice Lines", f"{m['total_records']:,}", Inches(6.9), Inches(1.1), box_color=DARK_BLUE)
    add_kpi(s, "Unique Vendors", f"{m.get('total_vendors', 0):,}", Inches(10.1), Inches(1.1), box_color=GREEN)

    add_kpi(s, "Direct Spend", f"EUR {m['direct_spend']/1e6:.1f}M ({m['direct_pct']:.0f}%)", Inches(0.5), Inches(2.1), box_color=GREEN)
    add_kpi(s, "Indirect Spend", f"EUR {m['indirect_spend']/1e6:.1f}M ({m['indirect_pct']:.0f}%)", Inches(3.7), Inches(2.1), box_color=ORANGE)
    add_kpi(s, "Debit Transactions", f"{m['debit_count']:,}", Inches(6.9), Inches(2.1), box_color=GREEN)
    add_kpi(s, "Credit Transactions", f"{m['credit_count']:,}", Inches(10.1), Inches(2.1), box_color=RED)

    bseg_txt = f"BSEG Coverage: GL Account {m['gl_pct']:.0f}% | Profit Center {m['prctr_pct']:.0f}% | Cost Center {m['kostl_pct']:.0f}%"
    add_insight_box(s, bseg_txt, Inches(0.5), Inches(3.1), Inches(6.0), Inches(0.7), LIGHT_BLUE, PPTX_COLORS['BSEG_BG'])

    # Indirect reasons table
    add_text(s, "Indirect Spend by Classification Reason", Inches(0.5), Inches(4.1), Inches(6), Inches(0.3), font_size=13, bold=True, color=SPAIN_RED)
    ind_reasons = m.get('indirect_reasons', pd.DataFrame())
    if len(ind_reasons) > 0:
        bd = [["Reason", "Lines", "EUR Amount", "% Indirect"]]
        for _, row in ind_reasons.head(5).iterrows():
            pct = row['EUR_Amount'] / m['indirect_spend'] * 100 if m['indirect_spend'] > 0 else 0
            bd.append([str(row['Reason'])[:35], f"{int(row['Lines']):,}", _f(row['EUR_Amount'], 'M1'), f"{pct:.1f}%"])
        add_table(s, bd, Inches(0.5), Inches(4.4), Inches(6.0), Inches(1.8), SPAIN_RED)

    # Top Vendors
    vs = m.get('vendor_spend', pd.DataFrame())
    if len(vs) > 0:
        add_text(s, "Top 10 Vendors by Spend", Inches(7.0), Inches(4.1), Inches(6), Inches(0.3), font_size=13, bold=True, color=DARK_BLUE)
        vd = [["Rank", "Vendor Name", "Gross EUR", "% Total"]]
        for i, (_, row) in enumerate(vs.head(10).iterrows(), 1):
            pct = row.get('Gross_EUR', row.get('Total_EUR', 0)) / m['total_gross'] * 100 if m['total_gross'] > 0 else 0
            vd.append([str(i), VN(row['Vendor'], 25), _f(row.get('Gross_EUR', row.get('Total_EUR', 0)), 'M1'), f"{pct:.1f}%"])
        add_table(s, vd, Inches(7.0), Inches(4.4), Inches(6.0), Inches(1.8), DARK_BLUE)

    add_footer(s, f"Source: SAP Final Comprehensive with BSEG + Vendor | BUKRS={entity_code} | {DATE}")
    return s
